"""Generate the Shiksha Sathi government pitch deck as a .pptx file.

Run from the repo root:
    python generate_pitch_deck.py

Outputs: pitch-deck.pptx  (open in PowerPoint, LibreOffice, or Google Slides)

Requires only python-pptx (already a project dependency via apps/api).
Install if needed:  pip install python-pptx
"""

import io
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ── Palette ───────────────────────────────────────────────────────────────────
ACCENT   = RGBColor(0xE5, 0x60, 0x10)   # Bihar saffron
GOLD     = RGBColor(0xF5, 0xA5, 0x2A)   # warm gold
NAVY     = RGBColor(0x0A, 0x16, 0x28)   # deep navy (slide bg)
NAVY_MID = RGBColor(0x13, 0x20, 0x40)   # card surface
TEXT     = RGBColor(0xF4, 0xED, 0xDF)   # warm off-white text
MUTED    = RGBColor(0x7A, 0x9A, 0xB8)   # slate-blue muted
GREEN    = RGBColor(0x50, 0xC8, 0x78)   # status green
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BORDER   = RGBColor(0x25, 0x3D, 0x5E)

# ── Dimensions ────────────────────────────────────────────────────────────────
W  = Inches(13.333)   # 16:9 widescreen
H  = Inches(7.5)
M  = Inches(0.6)      # margin

# EMU helpers
def emu(inches: float) -> int:
    return int(Inches(inches))


# ── Low-level helpers ─────────────────────────────────────────────────────────

def _bg(slide, color: RGBColor) -> None:
    """Fill the entire slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, x, y, w, h, fill: RGBColor, *, alpha: int = 255) -> None:
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()   # no border


def _textbox(
    slide, x, y, w, h, text: str, *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = TEXT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    italic: bool = False,
    wrap: bool = True,
) -> None:
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _label(slide, x, y, text: str) -> None:
    _textbox(slide, x, y, emu(4), emu(0.35), text,
             size=9, bold=True, color=ACCENT,
             align=PP_ALIGN.LEFT)


def _accent_bar(slide, height: float = 0.08) -> None:
    """Thin saffron stripe at the bottom of every slide."""
    bh = emu(height)
    _rect(slide, 0, int(H) - bh, int(W), bh, ACCENT)


def _brand_mark(slide) -> None:
    """'Shiksha Sathi' brand text in top-left corner."""
    _textbox(slide, emu(0.35), emu(0.18), emu(3.5), emu(0.35),
             "Shiksha Sathi  ·  शिक्षा सारथी",
             size=9, color=MUTED, bold=False)


def _card(slide, x, y, w, h) -> None:
    """Dark surface card with a border."""
    _rect(slide, x, y, w, h, NAVY_MID)
    # Thin border via a no-fill rectangle with a line
    border = slide.shapes.add_shape(1, x, y, w, h)
    border.fill.background()
    border.line.color.rgb = BORDER
    border.line.width = Pt(0.75)


# ── Blank slide builder ───────────────────────────────────────────────────────

def _blank(prs: Presentation):
    blank_layout = prs.slide_layouts[6]   # "Blank" layout
    slide = prs.slides.add_slide(blank_layout)
    _bg(slide, NAVY)
    return slide


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDES
# ═════════════════════════════════════════════════════════════════════════════

def slide_cover(prs):
    slide = _blank(prs)

    # Dot-grid texture: rows of small accent circles (approximate via tiny rects)
    dot_color = RGBColor(0x1C, 0x30, 0x50)
    step = emu(0.5)
    for col in range(0, int(W) + step, step * 2):
        for row in range(0, int(H) + step, step * 2):
            d = emu(0.04)
            s = slide.shapes.add_shape(9, col, row, d, d)  # 9 = oval
            s.fill.solid(); s.fill.fore_color.rgb = dot_color
            s.line.fill.background()

    # Central glow rectangle (soft)
    glow = slide.shapes.add_shape(9, emu(2.5), emu(1.2), emu(8.5), emu(5.2))
    glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(0x1A, 0x28, 0x45)
    glow.line.fill.background()

    # BSEB pill
    _rect(slide, emu(4.8), emu(1.3), emu(3.7), emu(0.35), NAVY_MID)
    _textbox(slide, emu(4.8), emu(1.3), emu(3.7), emu(0.35),
             "Bihar State Education Board  ·  BSEB Pilot",
             size=9, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # Hindi title
    _textbox(slide, emu(1.2), emu(1.9), emu(11), emu(1.8),
             "शिक्षा सारथी",
             size=64, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # English title
    _textbox(slide, emu(1.2), emu(3.4), emu(11), emu(1),
             "Shiksha Sathi",
             size=36, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

    # Hindi tagline
    _textbox(slide, emu(1.5), emu(4.3), emu(10.3), emu(0.6),
             "बिहार के सरकारी स्कूल शिक्षकों के लिए AI शिक्षण सहायक",
             size=16, color=GOLD, align=PP_ALIGN.CENTER)

    # English tagline
    _textbox(slide, emu(1.5), emu(4.85), emu(10.3), emu(0.5),
             "AI Teaching Companion for Bihar Government School Teachers  ·  Classes 6–10",
             size=13, color=MUTED, align=PP_ALIGN.CENTER)

    _accent_bar(slide)


def slide_problem(prs):
    slide = _blank(prs)
    _brand_mark(slide)

    _label(slide, M, emu(0.6), "THE CHALLENGE")
    _textbox(slide, M, emu(0.95), emu(7), emu(1),
             "Bihar's teachers are overburdened",
             size=32, bold=True, color=TEXT)
    _textbox(slide, M, emu(1.95), emu(9), emu(0.5),
             "हर रोज़ पाठ तैयारी में घंटों जाते हैं — डिजिटल उपकरण नहीं, हिंदी सामग्री नहीं",
             size=13, color=GOLD)

    # 3 stat cards
    stats = [
        ("70k+", "Government Schools", "sarkari schools across Bihar — most without teaching aids"),
        ("4L+",  "Teachers",           "preparing lessons manually, 3–4 hours per topic"),
        ("0",    "Hindi AI Tools",     "no AI tool built for BSEB curriculum or Hindi-medium teachers"),
    ]
    card_w = emu(3.8); card_h = emu(2.8)
    gap    = emu(0.25)
    start_x = M
    y = emu(2.8)
    for i, (num, title, desc) in enumerate(stats):
        x = int(start_x) + i * (card_w + gap)
        _card(slide, x, y, card_w, card_h)
        _textbox(slide, x + emu(0.2), y + emu(0.2), card_w - emu(0.4), emu(0.9),
                 num, size=40, bold=True, color=ACCENT)
        _textbox(slide, x + emu(0.2), y + emu(1.05), card_w - emu(0.4), emu(0.4),
                 title, size=13, bold=True, color=TEXT)
        _textbox(slide, x + emu(0.2), y + emu(1.5), card_w - emu(0.4), emu(1.0),
                 desc, size=11, color=MUTED)

    _textbox(slide, M, emu(6.3), emu(12), emu(0.4),
             "The result: inconsistent lesson quality, teacher burnout, and students left behind.",
             size=12, color=MUTED, align=PP_ALIGN.CENTER)
    _accent_bar(slide)


def slide_solution(prs):
    slide = _blank(prs)
    _brand_mark(slide)

    _label(slide, M, emu(0.6), "THE SOLUTION")
    _textbox(slide, M, emu(0.95), emu(8), emu(0.9),
             "One platform. Complete kit in minutes.",
             size=30, bold=True, color=TEXT)
    _textbox(slide, M, emu(1.8), emu(9), emu(0.4),
             "बोलें या टाइप करें — AI पूरा पाठ किट बनाएगा",
             size=13, color=GOLD)

    steps = [
        ("01", "अध्याय चुनें", "Select a Chapter",
         "Pick class, subject & chapter from the BSEB catalog — or speak it in Hindi via voice input"),
        ("02", "AI बनाए", "AI Generates",
         "6 resources stream in parallel — lesson plan, script, questions, worksheet, mind map, slides"),
        ("03", "पढ़ाएं!", "Teach!",
         "Download PPTX, print the worksheet, play the audio — or read the script aloud in class"),
    ]
    card_w = emu(3.8); card_h = emu(3.0)
    gap    = emu(0.25)
    y = emu(2.55)
    for i, (num, hi, en, desc) in enumerate(steps):
        x = int(M) + i * (card_w + gap)
        _card(slide, x, y, card_w, card_h)
        _textbox(slide, x + emu(0.2), y + emu(0.15), card_w - emu(0.4), emu(0.65),
                 num, size=28, bold=True, color=RGBColor(0x25, 0x3D, 0x5E))
        _textbox(slide, x + emu(0.2), y + emu(0.75), card_w - emu(0.4), emu(0.4),
                 hi, size=14, bold=True, color=GOLD)
        _textbox(slide, x + emu(0.2), y + emu(1.1), card_w - emu(0.4), emu(0.35),
                 en, size=13, bold=True, color=TEXT)
        _textbox(slide, x + emu(0.2), y + emu(1.5), card_w - emu(0.4), emu(1.2),
                 desc, size=11, color=MUTED)

    # Lightning bolt note
    _card(slide, M, emu(5.8), emu(12.1), emu(0.55))
    _textbox(slide, emu(0.8), emu(5.85), emu(11.5), emu(0.45),
             "⚡  Repeat lessons are instant — semantic caching serves a previously generated chapter in <1 second, no API cost.",
             size=11, color=TEXT)
    _accent_bar(slide)


def slide_kit(prs):
    slide = _blank(prs)
    _brand_mark(slide)

    _label(slide, M, emu(0.6), "TEACHING KIT — 7 RESOURCES GENERATED TOGETHER")
    _textbox(slide, M, emu(0.95), emu(9), emu(0.75),
             "Everything a teacher needs, already made",
             size=28, bold=True, color=TEXT)

    items = [
        ("📚", "पाठ योजना", "Lesson Plan",      "Structured 45-min plan with objectives & activities"),
        ("🗣️", "शिक्षण स्क्रिप्ट", "Teaching Script", "Word-for-word classroom script"),
        ("❓", "प्रश्नावली",  "Questions",        "MCQs, short-answer & HOTS with answer key"),
        ("📝", "वर्कशीट",    "Worksheet",        "Print-ready student worksheet (PDF)"),
        ("🗺️", "मन-मैप",    "Mind Map",         "Interactive visual concept map"),
        ("📊", "प्रेज़ेंटेशन", "Slides",           "5 / 10 / 15-slide PPTX — download & project"),
        ("🎧", "ऑडियो पाठ", "Audio Lesson",     "1, 3 and 5-minute spoken explanations (MP3)"),
    ]

    cols = 4
    card_w = emu(2.95); card_h = emu(1.9)
    gap_x  = emu(0.2);  gap_y  = emu(0.2)
    start_x = M;        start_y = emu(1.95)

    for idx, (icon, hi, en, desc) in enumerate(items):
        col = idx % cols
        row = idx // cols
        x = int(start_x) + col * (card_w + gap_x)
        y = int(start_y) + row * (card_h + gap_y)
        _card(slide, x, y, card_w, card_h)
        _textbox(slide, x + emu(0.15), y + emu(0.1), card_w - emu(0.3), emu(0.4),
                 icon + "  " + hi, size=12, bold=True, color=GOLD)
        _textbox(slide, x + emu(0.15), y + emu(0.52), card_w - emu(0.3), emu(0.3),
                 en, size=12, bold=True, color=TEXT)
        _textbox(slide, x + emu(0.15), y + emu(0.85), card_w - emu(0.3), emu(0.85),
                 desc, size=10, color=MUTED)

    _accent_bar(slide)


def slide_technology(prs):
    slide = _blank(prs)
    _brand_mark(slide)

    _label(slide, M, emu(0.6), "TECHNOLOGY")
    _textbox(slide, M, emu(0.95), emu(9), emu(0.8),
             "Built to scale across Bihar's 70,000 schools",
             size=28, bold=True, color=TEXT)

    tech = [
        ("🤖 Google Gemini AI",
         "State-of-the-art LLM generates BSEB-aligned content in Hindi, Hinglish and English"),
        ("⚡ Semantic Cache (pgvector)",
         "Cosine similarity search — teachers sharing the same chapter share the result, cutting AI cost to near-zero"),
        ("📶 Offline-Resilient PWA",
         "App shell cached locally — works on 2G/3G, shows a Hindi offline warning on poor connections"),
        ("🔐 Role-Based Access Control",
         "Teacher / School Admin / Super Admin roles — district officials monitor usage across their schools"),
        ("🎙️ Voice Input (OpenAI Whisper)",
         "Teachers speak in Hindi — transcribed & pre-filled automatically. No keyboard required."),
        ("📈 Analytics Dashboard",
         "Live cache hit rates by resource type — administrators see which topics are most requested"),
    ]

    card_w = emu(5.9); card_h = emu(1.45)
    gap_x  = emu(0.3); gap_y  = emu(0.25)
    start_x = M;       start_y = emu(2.1)

    for idx, (title, desc) in enumerate(tech):
        col = idx % 2
        row = idx // 2
        x = int(start_x) + col * (card_w + gap_x)
        y = int(start_y) + row * (card_h + gap_y)
        _card(slide, x, y, card_w, card_h)
        _textbox(slide, x + emu(0.2), y + emu(0.12), card_w - emu(0.4), emu(0.4),
                 title, size=12, bold=True, color=TEXT)
        _textbox(slide, x + emu(0.2), y + emu(0.55), card_w - emu(0.4), emu(0.75),
                 desc, size=10, color=MUTED)

    _accent_bar(slide)


def slide_bihar_first(prs):
    slide = _blank(prs)
    _brand_mark(slide)

    _label(slide, M, emu(0.6), "BIHAR-FIRST DESIGN")
    _textbox(slide, M, emu(0.95), emu(8), emu(0.8),
             "Not adapted for Bihar. Built for Bihar.",
             size=28, bold=True, color=TEXT)
    # accent underline on "Built for Bihar"
    _rect(slide, M, emu(1.75), emu(3.6), emu(0.05), ACCENT)

    features = [
        ("भा", "हिंदी और हिंग्लिश पहले",   "Hindi & Hinglish First",
         "UI, prompts and output all available in Hindi. Teachers generate content in the language they teach in."),
        ("📖", "BSEB पाठ्यक्रम",            "Full BSEB Curriculum Catalog",
         "All chapters for Classes 6–10 across Science, Math, Social Science, Hindi and English — real NCERT syllabus."),
        ("📱", "मोबाइल-पहले",               "Mobile-First Interface",
         "Teachers generate a full kit from their phone during free period. No laptop required."),
        ("🖨️", "प्रिंट-रेडी सामग्री",        "Print-Ready Materials",
         "Worksheets and PPTX slides formatted for A4 printing — for schools without projectors."),
    ]

    card_h = emu(1.1); gap = emu(0.2); start_y = emu(2.0)
    for i, (icon, hi, en, desc) in enumerate(features):
        y = int(start_y) + i * (card_h + gap)
        _card(slide, M, y, emu(12.1), card_h)
        # Icon circle
        _rect(slide, int(M) + emu(0.15), y + emu(0.2), emu(0.65), emu(0.65),
              RGBColor(0x1C, 0x35, 0x58))
        _textbox(slide, int(M) + emu(0.12), y + emu(0.12), emu(0.7), emu(0.7),
                 icon, size=14, align=PP_ALIGN.CENTER, bold=True, color=ACCENT)
        _textbox(slide, int(M) + emu(1.0), y + emu(0.08), emu(4), emu(0.35),
                 hi, size=11, bold=True, color=GOLD)
        _textbox(slide, int(M) + emu(1.0), y + emu(0.42), emu(3.0), emu(0.3),
                 en, size=11, bold=True, color=TEXT)
        _textbox(slide, int(M) + emu(4.5), y + emu(0.15), emu(7.4), emu(0.75),
                 desc, size=11, color=MUTED)

    _accent_bar(slide)


def slide_status(prs):
    slide = _blank(prs)
    _brand_mark(slide)

    _label(slide, M, emu(0.6), "PROJECT STATUS")
    _textbox(slide, M, emu(0.95), emu(8), emu(0.8),
             "8 of 9 phases complete.  Pilot-ready.",
             size=28, bold=True, color=TEXT)

    # 4 mini-stats
    mini = [
        ("80+", "Backend tests passing"),
        ("13",  "Frontend component tests"),
        ("E2E", "Playwright tested on mobile"),
        ("CI",  "ruff · mypy · tsc · eslint green"),
    ]
    card_w = emu(2.85); card_h = emu(1.1); gap = emu(0.2)
    y_mini = emu(2.0)
    for i, (num, lbl) in enumerate(mini):
        x = int(M) + i * (card_w + gap)
        _card(slide, x, y_mini, card_w, card_h)
        _textbox(slide, x + emu(0.2), y_mini + emu(0.05), card_w - emu(0.4), emu(0.5),
                 num, size=22, bold=True, color=ACCENT)
        _textbox(slide, x + emu(0.2), y_mini + emu(0.55), card_w - emu(0.4), emu(0.45),
                 lbl, size=10, color=MUTED)

    phases = [
        "Phase 1 — Monorepo Setup",
        "Phase 2 — Authentication",
        "Phase 3 — Database & Catalog",
        "Phase 4 — Kit Generation (all 7 nodes)",
        "Phase 5 — Audio Generation",
        "Phase 6 — PPT Hardening",
        "Phase 7 — Semantic Cache",
        "Phase 8 — Testing & RBAC",
    ]
    col_w = emu(5.85); ph_h = emu(0.45); ph_gap = emu(0.12)
    start_y = emu(3.35)
    for i, phase in enumerate(phases):
        col = i % 2; row = i // 2
        x = int(M) + col * (col_w + emu(0.4))
        y = int(start_y) + row * (ph_h + ph_gap)
        _card(slide, x, y, col_w, ph_h)
        # Green check
        _rect(slide, x + emu(0.12), y + emu(0.08), emu(0.3), emu(0.3), GREEN)
        _textbox(slide, x + emu(0.12), y + emu(0.04), emu(0.3), emu(0.3),
                 "✓", size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        _textbox(slide, x + emu(0.55), y + emu(0.06), col_w - emu(1.2), emu(0.35),
                 phase, size=10, color=TEXT)
        _textbox(slide, x + col_w - emu(0.9), y + emu(0.06), emu(0.8), emu(0.35),
                 "Done", size=9, color=GREEN, bold=True)

    _accent_bar(slide)


def slide_cta(prs):
    slide = _blank(prs)

    # Dot grid (lighter than cover)
    dot_color = RGBColor(0x1C, 0x30, 0x50)
    step = emu(0.5)
    for col in range(0, int(W) + step, step * 2):
        for row in range(0, int(H) + step, step * 2):
            d = emu(0.04)
            s = slide.shapes.add_shape(9, col, row, d, d)
            s.fill.solid(); s.fill.fore_color.rgb = dot_color
            s.line.fill.background()

    # Bottom glow
    glow = slide.shapes.add_shape(9, emu(1), emu(3.5), emu(11.3), emu(4.5))
    glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(0x14, 0x24, 0x3E)
    glow.line.fill.background()

    _brand_mark(slide)

    _textbox(slide, emu(1.2), emu(1.0), emu(11), emu(0.6),
             "Proposal",
             size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    _textbox(slide, emu(1.2), emu(1.5), emu(11), emu(1.4),
             "Launch a 100-school pilot",
             size=40, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

    _textbox(slide, emu(1.2), emu(2.85), emu(11), emu(0.55),
             "100 स्कूलों में पायलट — परिणाम देखें, फिर पूरे बिहार में विस्तार करें",
             size=14, color=GOLD, align=PP_ALIGN.CENTER)

    # 3 CTA steps
    cta_steps = [
        ("1", "Production Deploy",    "Vercel + Railway\nLive in 1 week"),
        ("2", "Pilot — 100 Schools",  "Onboard 100 teachers\ngather feedback"),
        ("3", "Scale to Bihar",       "Rollout to all\n70,000 schools"),
    ]
    card_w = emu(3.4); card_h = emu(1.6); gap = emu(0.3)
    start_x = emu(1.6); y_cta = emu(3.65)
    for i, (num, title, sub) in enumerate(cta_steps):
        x = int(start_x) + i * (card_w + gap)
        _card(slide, x, y_cta, card_w, card_h)
        # Circle number
        _rect(slide, x + int(card_w / 2) - emu(0.3), y_cta + emu(0.12),
              emu(0.6), emu(0.6), ACCENT)
        _textbox(slide, x + int(card_w / 2) - emu(0.3), y_cta + emu(0.1),
                 emu(0.6), emu(0.5),
                 num, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _textbox(slide, x + emu(0.15), y_cta + emu(0.78), card_w - emu(0.3), emu(0.38),
                 title, size=12, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        _textbox(slide, x + emu(0.15), y_cta + emu(1.15), card_w - emu(0.3), emu(0.35),
                 sub, size=10, color=MUTED, align=PP_ALIGN.CENTER)

    # URL badge
    _card(slide, emu(4.5), emu(5.55), emu(4.3), emu(0.45))
    _textbox(slide, emu(4.5), emu(5.55), emu(4.3), emu(0.45),
             "shiksha-sathi.vercel.app",
             size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    _textbox(slide, emu(1.2), emu(6.15), emu(11), emu(0.4),
             "Bihar State Education Board (BSEB)  ·  Government School Initiative  ·  Classes 6–10",
             size=10, color=MUTED, align=PP_ALIGN.CENTER)

    _accent_bar(slide)


# ═════════════════════════════════════════════════════════════════════════════
#  MAIN
# ═════════════════════════════════════════════════════════════════════════════

def build() -> bytes:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_cover(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_kit(prs)
    slide_technology(prs)
    slide_bihar_first(prs)
    slide_status(prs)
    slide_cta(prs)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


if __name__ == "__main__":
    out = Path("pitch-deck.pptx")
    out.write_bytes(build())
    print(f"Saved: {out.resolve()}")
    print(f"  {out.stat().st_size // 1024} KB  |  8 slides  |  Open in PowerPoint or Google Slides")
