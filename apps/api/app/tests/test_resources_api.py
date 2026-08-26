import io
import uuid

from fastapi.testclient import TestClient
from pptx import Presentation as read_presentation
from sqlalchemy.ext.asyncio import AsyncSession

from app.api.v1.deps import get_llm_provider
from app.main import app
from app.tests.fakes import FakeLLMProvider
from app.tests.test_teaching_kit_api import _create_profile, _seed_curriculum


def _generate_kit(client: TestClient, ids: tuple) -> tuple[str, list[dict]]:
    school_class, subject, chapter = ids
    app.dependency_overrides[get_llm_provider] = lambda: FakeLLMProvider()
    body = client.post(
        "/api/v1/teaching-kit/generate",
        json={
            "classId": str(school_class.id),
            "subjectId": str(subject.id),
            "chapterId": str(chapter.id),
            "language": "hi",
            "duration": "40",
        },
    ).json()
    request_id = body["requestId"]
    client.get(f"/api/v1/teaching-kit/{request_id}/stream")
    resources = client.get(f"/api/v1/teaching-kit/{request_id}").json()["resources"]
    return request_id, resources


def _resource_of(resources: list[dict], resource_type: str) -> dict:
    return next(r for r in resources if r["resourceType"] == resource_type)


async def test_read_resource_returns_its_content(
    client: TestClient, db_session: AsyncSession
) -> None:
    ids = await _seed_curriculum(db_session)
    _create_profile(client)
    _, resources = _generate_kit(client, ids)

    questions = _resource_of(resources, "questions")
    res = client.get(f"/api/v1/resources/{questions['id']}")

    assert res.status_code == 200
    assert res.json()["content"]["questions"][0]["type"] == "mcq"


async def test_read_resource_rejects_a_stranger(
    client: TestClient, db_session: AsyncSession
) -> None:
    ids = await _seed_curriculum(db_session)
    _create_profile(client)
    _, resources = _generate_kit(client, ids)

    # A resource is only reachable through the kit that produced it, so a
    # random id must 404 rather than leak that it exists.
    assert client.get(f"/api/v1/resources/{uuid.uuid4()}").status_code == 404


async def test_regenerate_produces_a_new_resource(
    client: TestClient, db_session: AsyncSession
) -> None:
    ids = await _seed_curriculum(db_session)
    _create_profile(client)
    _, resources = _generate_kit(client, ids)
    questions = _resource_of(resources, "questions")

    res = client.post(
        f"/api/v1/resources/{questions['id']}/regenerate",
        json={"params": {"difficulty": "advanced", "count": 5}},
    )

    assert res.status_code == 200
    body = res.json()
    assert body["id"] != questions["id"]
    # Different params derive a different cache key, so this is a real
    # generation rather than a hit on the original.
    assert body["cacheHit"] is False


async def test_regenerating_with_identical_params_hits_cache(
    client: TestClient, db_session: AsyncSession
) -> None:
    ids = await _seed_curriculum(db_session)
    _create_profile(client)
    _, resources = _generate_kit(client, ids)
    worksheet = _resource_of(resources, "worksheet")

    res = client.post(f"/api/v1/resources/{worksheet['id']}/regenerate", json={"params": {}})

    assert res.status_code == 200
    assert res.json()["cacheHit"] is True


async def test_audio_resource_can_be_regenerated(
    client: TestClient, db_session: AsyncSession
) -> None:
    ids = await _seed_curriculum(db_session)
    _create_profile(client)
    _, resources = _generate_kit(client, ids)
    audio = _resource_of(resources, "audio")

    res = client.post(f"/api/v1/resources/{audio['id']}/regenerate", json={"params": {}})

    assert res.status_code == 200
    content = res.json()["content"]
    assert set(content["variants"].keys()) == {"1", "3", "5"}


async def test_export_returns_a_real_pptx(client: TestClient, db_session: AsyncSession) -> None:
    ids = await _seed_curriculum(db_session)
    _create_profile(client)
    _, resources = _generate_kit(client, ids)
    presentation = _resource_of(resources, "presentation")

    res = client.get(f"/api/v1/resources/{presentation['id']}/export?format=pptx&version=5")

    assert res.status_code == 200
    assert res.headers["content-disposition"].startswith("attachment;")
    deck = read_presentation(io.BytesIO(res.content))
    assert len(deck.slides) == 5
    # Speaker notes are the point of generating a deck for someone teaching
    # from it rather than presenting it.
    assert deck.slides[0].notes_slide.notes_text_frame.text


async def test_export_rejects_a_non_presentation(
    client: TestClient, db_session: AsyncSession
) -> None:
    ids = await _seed_curriculum(db_session)
    _create_profile(client)
    _, resources = _generate_kit(client, ids)
    lesson_plan = _resource_of(resources, "lesson_plan")

    res = client.get(f"/api/v1/resources/{lesson_plan['id']}/export?format=pptx")

    assert res.status_code == 400
    assert res.json()["detail"] == "not_exportable_lesson_plan"


async def test_export_rejects_an_unsupported_format(
    client: TestClient, db_session: AsyncSession
) -> None:
    ids = await _seed_curriculum(db_session)
    _create_profile(client)
    _, resources = _generate_kit(client, ids)
    presentation = _resource_of(resources, "presentation")

    res = client.get(f"/api/v1/resources/{presentation['id']}/export?format=pdf")

    assert res.status_code == 400
    assert res.json()["detail"] == "unsupported_format"
