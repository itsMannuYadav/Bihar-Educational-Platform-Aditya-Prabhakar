import io

from pptx import Presentation as new_presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.presentation import Presentation
from pptx.util import Inches, Pt

# ── Colour palette ──────────────────────────────────────────────────────────
# The palette must project clearly from the back of a government-school
# classroom on a low-brightness wall (docs/07-roadmap.md Phase 6).
ACCENT = RGBColor(0xE5, 0x60, 0x10)    # warm saffron (readable on white)
DARK = RGBColor(0x1A, 0x1A, 0x2E)     # near-black navy for body text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ── Font sizes ───────────────────────────────────────────────────────────────
TITLE_SLIDE_TITLE_PT = 44
BODY_SLIDE_TITLE_PT = 32
BODY_TEXT_PT = 24
SPEAKER_NOTES_PT = 12  # not enlarged — only the teacher reads these

# ── Layout indices ───────────────────────────────────────────────────────────
_TITLE_LAYOUT = 0
_TITLE_AND_CONTENT_LAYOUT = 1

# Accent bar geometry (a thin strip across the bottom for visual grounding)
_BAR_HEIGHT = Inches(0.12)
_SLIDE_W = Inches(13.333)
_SLIDE_H = Inches(7.5)


def _accent_bar(slide) -> None:  # type: ignore[type-arg]
    """Add a thin saffron stripe at the very bottom of a slide."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0,
        int(_SLIDE_H - _BAR_HEIGHT),
        int(_SLIDE_W),
        int(_BAR_HEIGHT),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    # Remove border
    sp_tree = bar._element
    sp_pr = sp_tree.find(qn("p:spPr"))
    if sp_pr is not None:
        ln = sp_pr.find(qn("a:ln"))
        if ln is not None:
            sp_pr.remove(ln)
    bar.line.color.rgb = ACCENT  # fallback: match fill so no border shows


def _set_text_color(frame, color: RGBColor) -> None:
    for para in frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = color


class NativePptxProvider:
    """Renders the graph's provider-agnostic slide list into a real .pptx.

    Text is written as text, never as rasterized images, so PowerPoint (or
    LibreOffice, or Google Slides) does the script shaping — which is what
    makes Hindi decks render correctly without shipping fonts.

    Phase 6 design pass: saffron accent bars, high-contrast navy body text,
    large fonts sized for back-of-classroom projection.
    """

    extension = "pptx"
    media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    def render(self, *, title: str, slides: list[dict]) -> bytes:
        presentation = new_presentation()
        presentation.slide_width = _SLIDE_W
        presentation.slide_height = _SLIDE_H

        for index, slide_data in enumerate(slides):
            self._add_slide(presentation, slide_data, is_first=index == 0, deck_title=title)

        buffer = io.BytesIO()
        presentation.save(buffer)
        return buffer.getvalue()

    def _add_slide(
        self, presentation: Presentation, slide_data: dict, *, is_first: bool, deck_title: str
    ) -> None:
        layout_name = slide_data.get("layout", "bullets")
        heading = slide_data.get("title") or (deck_title if is_first else "")
        body: list[str] = slide_data.get("body") or []

        if layout_name == "title" or is_first:
            slide = presentation.slides.add_slide(presentation.slide_layouts[_TITLE_LAYOUT])

            title_tf = slide.shapes.title.text_frame
            title_tf.clear()
            p = title_tf.paragraphs[0]
            p.text = heading
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(TITLE_SLIDE_TITLE_PT)
            run.font.color.rgb = ACCENT
            run.font.bold = True

            if len(slide.placeholders) > 1 and body:
                subtitle = slide.placeholders[1]
                subtitle.text = body[0]
                subtitle.text_frame.paragraphs[0].font.size = Pt(BODY_TEXT_PT)
                _set_text_color(subtitle.text_frame, DARK)
        else:
            slide = presentation.slides.add_slide(
                presentation.slide_layouts[_TITLE_AND_CONTENT_LAYOUT]
            )

            title_tf = slide.shapes.title.text_frame
            title_tf.clear()
            p = title_tf.paragraphs[0]
            p.text = heading
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(BODY_SLIDE_TITLE_PT)
            run.font.color.rgb = DARK
            run.font.bold = True

            text_frame = slide.placeholders[1].text_frame
            text_frame.word_wrap = True
            for i, line in enumerate(body or [""]):
                para = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                para.text = line
                para.line_spacing = Pt(BODY_TEXT_PT * 1.4)
                for run in para.runs:
                    run.font.size = Pt(BODY_TEXT_PT)
                    run.font.color.rgb = DARK

        _accent_bar(slide)

        # Speaker notes carry what the teacher actually says — the whole point
        # of generating a deck for someone teaching from it, not presenting it.
        if notes := slide_data.get("speaker_notes"):
            notes_tf = slide.notes_slide.notes_text_frame
            notes_tf.text = notes
            for para in notes_tf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(SPEAKER_NOTES_PT)
